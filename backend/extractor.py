"""Lumina OCR & Document Extraction Engine.

Extracts text from any file format:
  - PDF  (text layer + OCR fallback for scanned pages)
  - DOCX / PPTX / XLSX / CSV / TXT / MD
  - Images (PNG/JPG/JPEG/WEBP/BMP/TIFF) via RapidOCR (no external binaries)
  - HTML / RTF

Strategy per file:
  1. Detect real type from magic bytes (not just extension).
  2. Prefer fast, lossless embedded-text extraction (pdf text layer, docx xml, ...).
  3. Fall back to OCR when a page/image has no (or too little) extractable text.
"""

from __future__ import annotations

import csv
import io
import os
import re
import tempfile
import threading
import zipfile
from dataclasses import dataclass, field
from typing import List, Optional, Tuple

# ---------------------------------------------------------------------------
# Lazy imports (heavy modules load only when actually needed)
# ---------------------------------------------------------------------------

_PDF_OK = True
try:
    import fitz  # PyMuPDF
except Exception:  # pragma: no cover
    _PDF_OK = False

_DOCX_OK = True
try:
    import docx
except Exception:  # pragma: no cover
    _DOCX_OK = False

_PPTX_OK = True
try:
    from pptx import Presentation
except Exception:  # pragma: no cover
    _PPTX_OK = False

_XLSX_OK = True
try:
    from openpyxl import load_workbook
except Exception:  # pragma: no cover
    _XLSX_OK = False

_OCR_OK = True
try:
    from rapidocr_onnxruntime import RapidOCR
except Exception:  # pragma: no cover
    _OCR_OK = False

_HTML_OK = True
try:
    from bs4 import BeautifulSoup
except Exception:  # pragma: no cover
    _HTML_OK = False

# ---------------------------------------------------------------------------
# Result model
# ---------------------------------------------------------------------------


@dataclass
class ExtractResult:
    text: str = ""
    method: str = "none"            # pdf-text | pdf-ocr | docx | pptx | xlsx | csv | image-ocr | plain | html | zip
    pages: int = 0
    chars: int = 0
    ocrPages: int = 0
    warnings: List[str] = field(default_factory=list)
    language: str = "en"

    def to_dict(self) -> dict:
        return {
            "text": self.text,
            "method": self.method,
            "pages": self.pages,
            "chars": self.chars,
            "ocrPages": self.ocrPages,
            "warnings": self.warnings,
            "language": self.language,
        }


# ---------------------------------------------------------------------------
# OCR engine (singleton, lazy init, thread-safe)
# ---------------------------------------------------------------------------

_ocr: Optional[RapidOCR] = None
_ocr_lock = threading.Lock()


_OCR_FAILED = object()


def _get_ocr():
    global _ocr
    if not _OCR_OK:
        return None
    if _ocr is None:
        with _ocr_lock:
            if _ocr is None:
                try:
                    _ocr = RapidOCR()
                except Exception:
                    _ocr = _OCR_FAILED
    return None if _ocr is _OCR_FAILED else _ocr


def _ocr_image_bytes(data: bytes) -> Tuple[str, int]:
    """Run OCR on raw image bytes. Returns (text, block_count)."""
    ocr = _get_ocr()
    if ocr is None:
        return "", 0
    try:
        import numpy as np
        from PIL import Image

        img = Image.open(io.BytesIO(data)).convert("RGB")
        result, _ = ocr(np.array(img))
        if not result:
            return "", 0
        lines: List[Tuple[float, float, str]] = []
        for box, text, score in result:
            if text is None or not str(text).strip():
                continue
            y = float(box[0][1])
            x = float(box[0][0])
            lines.append((y, x, str(text)))
        lines.sort(key=lambda t: (round(t[0] / 12), t[1]))
        return "\n".join(t[2] for t in lines), len(lines)
    except Exception:
        return "", 0


# ---------------------------------------------------------------------------
# Type detection (magic bytes first)
# ---------------------------------------------------------------------------

MAGIC_SIGNATURES = [
    (b"%PDF", "pdf"),
    (b"\x50\x4b\x03\x04", "zip"),      # docx/pptx/xlsx/open office
    (b"\x89PNG\r\n\x1a\n", "png"),
    (b"\xff\xd8\xff", "jpg"),
    (b"GIF87a", "gif"),
    (b"GIF89a", "gif"),
    (b"BM", "bmp"),
    (b"RIFF", "webp"),
    (b"II*\x00", "tiff"),
    (b"MM\x00*", "tiff"),
    (b"\x7fELF", "elf"),
]

EXT_TO_TYPE = {
    ".pdf": "pdf",
    ".docx": "docx",
    ".doc": "doc",
    ".pptx": "pptx",
    ".ppt": "ppt",
    ".xlsx": "xlsx",
    ".xls": "xls",
    ".csv": "csv",
    ".tsv": "tsv",
    ".txt": "txt",
    ".md": "md",
    ".markdown": "md",
    ".html": "html",
    ".htm": "html",
    ".xml": "xml",
    ".json": "json",
    ".rtf": "rtf",
    ".png": "png",
    ".jpg": "jpg",
    ".jpeg": "jpg",
    ".webp": "webp",
    ".bmp": "bmp",
    ".tiff": "tiff",
    ".tif": "tiff",
    ".gif": "gif",
    ".epub": "epub",
    ".odt": "odt",
    ".ods": "ods",
    ".odp": "odp",
}

IMAGE_TYPES = {"png", "jpg", "gif", "bmp", "webp", "tiff", "jpeg"}
OFFICE_TYPES = {"docx", "pptx", "xlsx"}


def detect_type(filename: str, data: bytes) -> str:
    head = data[:16] if data else b""
    for sig, ftype in MAGIC_SIGNATURES:
        if head.startswith(sig):
            if ftype == "zip":
                inner = _probe_zip(data)
                if inner:
                    return inner
            return ftype
    ext = os.path.splitext(filename or "")[1].lower()
    if ext == ".jpg":
        return "jpg"
    return EXT_TO_TYPE.get(ext, "txt")


def _probe_zip(data: bytes) -> Optional[str]:
    """Look inside a zip to decide docx / pptx / xlsx / odt."""
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as z:
            names = z.namelist()
        if any(n.startswith("word/") for n in names):
            return "docx"
        if any(n.startswith("ppt/") for n in names):
            return "pptx"
        if any(n.startswith("xl/") for n in names):
            return "xlsx"
        if any(n.startswith("content.xml") and "office" in str(names) for n in names):
            return "odt"
    except Exception:
        pass
    return None


# ---------------------------------------------------------------------------
# Plain text helpers
# ---------------------------------------------------------------------------

def _clean(text: str) -> str:
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"[ \t]+\n", "\n", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _extract_plain(data: bytes) -> str:
    for enc in ("utf-8-sig", "utf-16", "latin-1"):
        try:
            return data.decode(enc)
        except Exception:
            continue
    return data.decode("utf-8", errors="replace")


def _extract_rtf(data: bytes) -> str:
    text = data.decode("latin-1", errors="replace")
    text = re.sub(r"\\[a-z]+-?\d* ?", "", text)
    text = text.replace("{", "").replace("}", "").replace("\\par", "\n").replace("\\tab", "\t")
    return _clean(text)


# ---------------------------------------------------------------------------
# CSV / TSV
# ---------------------------------------------------------------------------

def _extract_csv(data: bytes, sep: str = ",") -> str:
    try:
        text = data.decode("utf-8-sig")
    except Exception:
        text = data.decode("latin-1", errors="replace")
    rows = list(csv.reader(io.StringIO(text), delimiter=sep))
    out = []
    for row in rows:
        out.append(" | ".join(c.strip() for c in row))
    return _clean("\n".join(out))


# ---------------------------------------------------------------------------
# HTML
# ---------------------------------------------------------------------------

def _extract_html(data: bytes) -> str:
    if not _HTML_OK:
        return _extract_plain(data)
    soup = BeautifulSoup(data.decode("utf-8", errors="replace"), "html.parser")
    for tag in soup(["script", "style", "noscript", "svg"]):
        tag.decompose()
    for tag in soup.find_all(["h1", "h2", "h3", "h4", "h5", "h6"]):
        tag.append("\n" + "#" * int(tag.name[1]) + " ")
    for tag in soup.find_all(["li"]):
        tag.append("\n- ")
    for tag in soup.find_all(["tr"]):
        tag.append("\n")
    for tag in soup.find_all(["p", "div", "br", "table", "blockquote"]):
        tag.append("\n")
    text = soup.get_text(separator=" ")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r" *\n *", "\n", text)
    return _clean(text)


# ---------------------------------------------------------------------------
# Office formats (docx / pptx / xlsx)
# ---------------------------------------------------------------------------

def _extract_docx(data: bytes) -> str:
    if not _DOCX_OK:
        return _extract_plain(data)
    doc = docx.Document(io.BytesIO(data))
    parts: List[str] = []

    for para in doc.paragraphs:
        t = para.text.strip()
        if t:
            parts.append(t)
    for table in doc.tables:
        parts.append("")
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            parts.append(" | ".join(cells))
    # Headers / footers
    for section in doc.sections:
        for para in section.header.paragraphs:
            if para.text.strip():
                parts.append(f"[Header] {para.text.strip()}")
        for para in section.footer.paragraphs:
            if para.text.strip():
                parts.append(f"[Footer] {para.text.strip()}")
    return _clean("\n".join(parts))


def _iter_pptx_shapes(shapes) -> List[str]:
    out: List[str] = []
    for shape in shapes:
        if shape.shape_type == 19:  # GROUP
            try:
                out.extend(_iter_pptx_shapes(shape.shapes))
            except Exception:
                pass
            continue
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = "".join(run.text for run in para.runs).strip()
                if t:
                    out.append(t)
        if getattr(shape, "has_table", False) and shape.has_table:
            out.append("")
            for row in shape.table.rows:
                out.append(" | ".join(cell.text.strip() for cell in row.cells))
        if getattr(shape, "has_chart", False):
            try:
                chart = shape.chart
                plot = chart.plots[0]
                cats = [str(c) for c in plot.categories]
                out.append(" | ".join(cats))
                for series in plot.series:
                    out.append(f"{series.name}: " + ", ".join(str(v) for v in series.values))
            except Exception:
                pass
    return out


def _extract_pptx(data: bytes) -> str:
    if not _PPTX_OK:
        return _extract_plain(data)
    prs = Presentation(io.BytesIO(data))
    parts: List[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        parts.append(f"\n--- Slide {idx} ---")
        parts.extend(_iter_pptx_shapes(slide.shapes))
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"[Speaker notes] {notes}")
    return _clean("\n".join(parts))


def _extract_xlsx(data: bytes) -> str:
    if not _XLSX_OK:
        return _extract_plain(data)
    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    parts: List[str] = []
    for ws in wb.worksheets:
        parts.append(f"\n=== Sheet: {ws.title} ===")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c).strip() if c is not None else "" for c in row]
            if any(cells):
                parts.append(" | ".join(cells))
    wb.close()
    return _clean("\n".join(parts))


# ---------------------------------------------------------------------------
# PDF (text layer + OCR fallback)
# ---------------------------------------------------------------------------

MIN_TEXT_CHARS_PER_PAGE = 12  # below this, treat page as scanned -> OCR
MAX_OCR_PAGES = 60            # safety cap
OCR_ZOOM = 2.5                # ~180 dpi


def _extract_pdf(data: bytes, ocr: bool = True) -> ExtractResult:
    if not _PDF_OK:
        return ExtractResult(text="", method="pdf-unavailable", warnings=["PyMuPDF not installed"])
    res = ExtractResult(method="pdf-text")
    doc = fitz.open(stream=data, filetype="pdf")
    res.pages = doc.page_count
    parts: List[str] = []
    total_text = 0
    for pno in range(doc.page_count):
        page = doc[pno]
        text = page.get_text("text") or ""
        total_text += len(text.strip())
        parts.append(f"\n--- Page {pno + 1} ---\n{text.strip()}")
    if ocr and _OCR_OK:
        for pno in range(doc.page_count):
            if res.ocrPages >= MAX_OCR_PAGES:
                res.warnings.append(f"OCR capped at {MAX_OCR_PAGES} pages")
                break
            page = doc[pno]
            text = page.get_text("text") or ""
            if len(text.strip()) >= MIN_TEXT_CHARS_PER_PAGE:
                continue
            try:
                pix = page.get_pixmap(matrix=fitz.Matrix(OCR_ZOOM, OCR_ZOOM))
                img_bytes = pix.tobytes("png")
            except Exception:
                continue
            ocr_text, _ = _ocr_image_bytes(img_bytes)
            if ocr_text.strip():
                res.ocrPages += 1
                parts.append(f"\n--- Page {pno + 1} (OCR) ---\n{ocr_text}")
    doc.close()
    res.text = _clean("\n".join(parts))
    res.chars = len(res.text)
    if res.ocrPages > 0:
        res.method = "pdf-text+ocr"
    return res


# ---------------------------------------------------------------------------
# Images
# ---------------------------------------------------------------------------

def _extract_image(data: bytes, ftype: str) -> ExtractResult:
    if not _OCR_OK:
        return ExtractResult(text="", method="image-ocr", warnings=["OCR engine not installed"])
    res = ExtractResult(method="image-ocr", pages=1)
    text, blocks = _ocr_image_bytes(data)
    res.text = _clean(text)
    res.chars = len(res.text)
    if not res.text:
        res.warnings.append("No text detected in image")
    return res


# ---------------------------------------------------------------------------
# Zips without office content (epub, odt) — conservative fallback
# ---------------------------------------------------------------------------

def _extract_odt(data: bytes) -> str:
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as z:
            xml = z.read("content.xml").decode("utf-8", errors="replace")
        if _HTML_OK:
            soup = BeautifulSoup(xml, "xml") if "lxml" in os.sys.modules else BeautifulSoup(xml, "html.parser")
            for tag in soup.find_all(["text:p", "text:h"]):
                tag.append("\n")
            for tag in soup.find_all("table:table-cell"):
                tag.append(" | ")
            return _clean(soup.get_text(separator=" "))
        return _clean(re.sub(r"<[^>]+>", " ", xml))
    except Exception:
        return ""


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

def extract_bytes(filename: str, data: bytes, ocr: bool = True) -> ExtractResult:
    ftype = detect_type(filename, data)
    try:
        if ftype == "pdf":
            return _extract_pdf(data, ocr=ocr)
        if ftype == "docx":
            res = ExtractResult(text=_extract_docx(data), method="docx", pages=1)
        elif ftype == "pptx":
            res = ExtractResult(text=_extract_pptx(data), method="pptx", pages=0)
        elif ftype == "xlsx":
            res = ExtractResult(text=_extract_xlsx(data), method="xlsx", pages=0)
        elif ftype in ("csv", "tsv"):
            res = ExtractResult(text=_extract_csv(data, sep="\t" if ftype == "tsv" else ","), method="csv", pages=0)
        elif ftype in IMAGE_TYPES:
            return _extract_image(data, ftype)
        elif ftype in ("html", "htm", "xml"):
            res = ExtractResult(text=_extract_html(data), method="html", pages=1)
        elif ftype == "rtf":
            res = ExtractResult(text=_extract_rtf(data), method="rtf", pages=1)
        elif ftype in ("odt", "odp", "ods"):
            res = ExtractResult(text=_extract_odt(data), method="office-open", pages=0)
        else:
            res = ExtractResult(text=_extract_plain(data), method="plain", pages=1)
        res.chars = len(res.text)
        if not res.text and ftype in OFFICE_TYPES:
            res.warnings.append("No readable text found (file may be image-only)")
        return res
    except Exception as exc:  # pragma: no cover
        return ExtractResult(text="", method="error", warnings=[f"Extraction failed: {exc}"])


def extract_file(path: str, ocr: bool = True) -> ExtractResult:
    with open(path, "rb") as fh:
        data = fh.read()
    return extract_bytes(os.path.basename(path), data, ocr=ocr)


def extract_tempfile_bytes(suffix: str, data: bytes) -> ExtractResult:
    """Workaround for engines that need a real file on disk (rare)."""
    fd, path = tempfile.mkstemp(suffix=suffix)
    try:
        with os.fdopen(fd, "wb") as fh:
            fh.write(data)
        return extract_file(path)
    finally:
        try:
            os.remove(path)
        except Exception:
            pass


# ---------------------------------------------------------------------------
# URL fetch + extract convenience
# ---------------------------------------------------------------------------

def extract_url(url: str, ocr: bool = True, timeout: float = 30.0) -> ExtractResult:
    import urllib.request

    req = urllib.request.Request(url, headers={"User-Agent": "Mozilla/5.0 (Lumina-Extractor)"})
    with urllib.request.urlopen(req, timeout=timeout) as resp:
        data = resp.read()
        ctype = resp.headers.get("Content-Type", "")
    name = url.split("/")[-1].split("?")[0]
    if not name:
        name = "remote-file" + (".pdf" if "pdf" in ctype else ".html")
    return extract_bytes(name, data, ocr=ocr)
